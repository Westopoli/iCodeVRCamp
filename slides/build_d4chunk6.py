"""Build D4Chunk6.pptx — single reference slide: TODO #6 state-machine code shapes.

Four cards (2 × 2 grid), one per match branch (6a–6d), showing the code
shape kids fill inside each arm of the state machine in player.gd.

Mirrors the D2FC / D3FC / D4FC methodology: red label bar, dark code bg,
syntax-highlighted GDScript, ← mirror note.

Usage:
    cd slides
    .venv/bin/python build_d4chunk6.py   -> out/D4Chunk6.pptx
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import theme
import master

# ── data ──────────────────────────────────────────────────────────────────────
# Code shapes are the verbatim fill for each #@todo region inside the
# `match state:` dispatcher in player.gd KID CHUNK #6.

EXAMPLES = [
    {
        "label": '#6a — "idle" state exits',
        "mirror": "Chunk #6  ·  inside  match state:",
        "code": [
            '"idle":',
            "    if get_move_direction() != 0:",
            '        set_state("walk")',
            '    if get_input_just_pressed("jump") and is_on_floor():',
            "        velocity.y = -jump_impulse",
            '        set_state("jump")',
        ],
    },
    {
        "label": '#6b — "walk" state exits',
        "mirror": "Chunk #6  ·  inside  match state:",
        "code": [
            '"walk":',
            "    if get_move_direction() == 0:",
            '        set_state("idle")',
            '    if get_input_just_pressed("jump") and is_on_floor():',
            "        velocity.y = -jump_impulse",
            '        set_state("jump")',
        ],
    },
    {
        "label": '#6c — "jump" state exit',
        "mirror": "Chunk #6  ·  inside  match state:",
        "code": [
            '"jump":',
            "    if velocity.y > 0:",
            '        set_state("fall")',
        ],
    },
    {
        "label": '#6d — "fall" state exit',
        "mirror": "Chunk #6  ·  inside  match state:",
        "code": [
            '"fall":',
            "    if is_on_floor():",
            '        set_state("idle")',
        ],
    },
]

# ── layout ────────────────────────────────────────────────────────────────────

COLS = 2
ROWS = 2
MARGIN_H = Inches(0.35)
MARGIN_TOP_BODY = Inches(0.18)
GAP_COL = Inches(0.22)
GAP_ROW = Inches(0.22)
MARGIN_BOTTOM = Inches(0.18)

BODY_W = theme.SLIDE_WIDTH - MARGIN_H * 2
BODY_H = theme.SLIDE_HEIGHT - theme.HEADER_HEIGHT - MARGIN_TOP_BODY - MARGIN_BOTTOM

CELL_W = (BODY_W - GAP_COL * (COLS - 1)) / COLS
CELL_H = (BODY_H - GAP_ROW * (ROWS - 1)) / ROWS

LABEL_H = Inches(0.40)
MIRROR_H = Inches(0.28)
PAD = Inches(0.12)
CODE_SIZE = Pt(12)
LABEL_SIZE = Pt(12)

# ── tokeniser (matches D4FC) ───────────────────────────────────────────────────

_KEYWORDS = {
    "for", "in", "range", "var", "const", "func", "return",
    "if", "else", "elif", "and", "or", "not", "pass",
    "true", "false", "null",
}
_FLOW_KW = {"while", "if", "else", "elif", "match"}
_TYPES = {"bool", "int", "float", "String", "void", "Vector2i", "Vector2", "Array", "Node", "Color"}


def tokenize(line: str):
    if not line.strip():
        return [(" ", theme.CODE_TEXT)]
    result = []
    stripped = line.lstrip(" ")
    indent = len(line) - len(stripped)
    if indent:
        result.append((" " * indent, theme.CODE_TEXT))
    line = stripped
    i = 0
    while i < len(line):
        ch = line[i]
        if ch == "#":
            result.append((line[i:], theme.SYNTAX_COMMENT))
            break
        if ch in ('"', "'"):
            q = ch
            j = i + 1
            while j < len(line) and line[j] != q:
                j += 1
            j += 1
            result.append((line[i:j], theme.SYNTAX_STRING))
            i = j
            continue
        if ch.isdigit():
            j = i
            while j < len(line) and (line[j].isdigit() or line[j] == "."):
                j += 1
            result.append((line[i:j], theme.SYNTAX_NUMBER))
            i = j
            continue
        if ch.isalpha() or ch == "_":
            j = i
            while j < len(line) and (line[j].isalnum() or line[j] == "_"):
                j += 1
            word = line[i:j]
            if word in _TYPES:
                color = theme.SYNTAX_TYPE
            elif word in _FLOW_KW:
                color = theme.SYNTAX_KEYWORD_FLOW
            elif word in _KEYWORDS:
                color = theme.SYNTAX_KEYWORD
            elif j < len(line) and line[j] == "(":
                color = theme.SYNTAX_FUNCTION
            else:
                color = theme.CODE_TEXT
            result.append((word, color))
            i = j
            continue
        two = line[i:i+2]
        if two in ("->", ":=", "+=", "-=", "*=", "==", "!=", "<=", ">="):
            result.append((two, theme.SYNTAX_OPERATOR))
            i += 2
            continue
        result.append((ch, theme.CODE_TEXT))
        i += 1
    return result


# ── builder ───────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = theme.SLIDE_WIDTH
    prs.slide_height = theme.SLIDE_HEIGHT

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    master.apply_master(slide, day_number=4, page_number=None)

    # Title in the header bar
    title_box = slide.shapes.add_textbox(
        Inches(3.0), Inches(0.18), Inches(7.5), Inches(0.55),
    )
    tf = title_box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "TODO #6 — State Machine  (player.gd)"
    run.font.name = theme.FONT_HEADING
    run.font.size = Pt(19)
    run.font.bold = True
    run.font.color.rgb = theme.BG_WHITE

    for idx, ex in enumerate(EXAMPLES):
        col = idx % COLS
        row = idx // COLS

        left = MARGIN_H + col * (CELL_W + GAP_COL)
        top = theme.HEADER_HEIGHT + MARGIN_TOP_BODY + row * (CELL_H + GAP_ROW)

        # Dark card background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, CELL_W, CELL_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.CODE_BG
        bg.line.fill.background()
        bg.shadow.inherit = False

        # Red label bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, CELL_W, LABEL_H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme.ICODE_RED
        bar.line.fill.background()
        bar.shadow.inherit = False

        # Label text (bold number + regular description)
        lbl = slide.shapes.add_textbox(left + PAD, top, CELL_W - PAD * 2, LABEL_H)
        tf = lbl.text_frame
        tf.margin_top = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        fc_num, _, label_rest = ex["label"].partition(" — ")
        run = p.add_run()
        run.text = fc_num + "  "
        run.font.name = theme.FONT_HEADING
        run.font.size = LABEL_SIZE
        run.font.bold = True
        run.font.color.rgb = theme.BG_WHITE
        run2 = p.add_run()
        run2.text = label_rest
        run2.font.name = theme.FONT_HEADING
        run2.font.size = LABEL_SIZE
        run2.font.bold = False
        run2.font.color.rgb = theme.BG_WHITE

        # Mirror note
        mirror_top = top + LABEL_H
        mirror_box = slide.shapes.add_textbox(
            left + PAD, mirror_top, CELL_W - PAD * 2, MIRROR_H,
        )
        tf = mirror_box.text_frame
        tf.margin_top = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"← {ex['mirror']}"
        run.font.name = theme.FONT_BODY
        run.font.size = Pt(10)
        run.font.italic = True
        run.font.color.rgb = theme.GREY_MID

        # Code block
        code_top = top + LABEL_H + MIRROR_H + PAD
        code_h = CELL_H - LABEL_H - MIRROR_H - PAD * 2
        code_box = slide.shapes.add_textbox(
            left + PAD, code_top, CELL_W - PAD * 2, code_h
        )
        tf = code_box.text_frame
        tf.margin_top = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.word_wrap = False

        for li, line in enumerate(ex["code"]):
            para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            for text, color in tokenize(line):
                run = para.add_run()
                run.text = text
                run.font.name = theme.FONT_MONO
                run.font.size = CODE_SIZE
                run.font.color.rgb = color

    out = Path(__file__).parent / "out" / "D4Chunk6.pptx"
    prs.save(str(out))
    print(f"Wrote {out} — 1 slide")


if __name__ == "__main__":
    build()
