"""Extract hand-positioned red-overlay boxes from an existing built deck and save
their geometry keyed by stable slide-ID, so future rebuilds restore them.

The build draws each red overlay as a rectangle: no fill, red 4pt outline, no
text. (The dashed MISSING-placeholder also uses a red line but is 2pt + grey
fill + text, so it's filtered out.)

Matching old slide -> slide-ID:
  1. Hash the screenshot embedded on the slide; map it to a file on disk; map
     that filename to the blueprint slide that references it (via the same
     resolution build_day uses). Robust + order-independent.
  2. If a screenshot is shared by 2 slides (where + do-it), disambiguate by the
     slide's title text.
  3. Slides whose box sits on a placeholder (no embedded image) fall back to a
     unique-title match, else are reported as unmatched.

Usage:
    cd slides
    python extract_overlays.py 1            # reads out/Day1.pptx
    python extract_overlays.py 1 SomeDeck.pptx
"""

import hashlib
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

import build_day as bd
import theme

HERE = Path(__file__).parent
RED = (theme.OVERLAY_RED[0], theme.OVERLAY_RED[1], theme.OVERLAY_RED[2])
MIN_OVERLAY_LINE = Pt(3)  # real overlay is 4pt; placeholder border is 2pt


def _is_overlay(shape):
    """True if shape looks like a hand-editable red overlay box."""
    try:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return False  # placeholders carry "[ MISSING: ... ]" text
    except Exception:
        pass
    try:
        line = shape.line
        if line.color.rgb != theme.OVERLAY_RED:
            return False
        if line.width is None or int(line.width) < int(MIN_OVERLAY_LINE):
            return False
    except Exception:
        return False
    try:
        # overlay has no fill (background); placeholder has a solid grey fill
        if shape.fill.type is not None and shape.fill.type != 5:  # 5 = MSO_FILL.BACKGROUND
            return False
    except Exception:
        pass
    return True


import re as _re
_CHROME = _re.compile(r"^(icode|empowering future innovators|day\s*\d+|\d+)$", _re.I)


def _slide_title(slide):
    """First real content title, skipping master-frame chrome (logo wordmark,
    tagline, 'DAY N' label, page number)."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            first = sh.text_frame.text.strip().splitlines()[0].strip()
            if first and not _CHROME.match(first):
                return first
    return ""


def _slide_image_hashes(slide):
    out = []
    for sh in slide.shapes:
        if sh.shape_type == 13:  # PICTURE
            try:
                out.append(hashlib.md5(sh.image.blob).hexdigest())
            except Exception:
                pass
    return out


def _disk_hash_index():
    """md5 -> filename for every screenshot on disk."""
    idx = {}
    for d in (HERE / "screenshots").rglob("*.png"):
        try:
            idx[hashlib.md5(d.read_bytes()).hexdigest()] = d.name
        except Exception:
            pass
    return idx


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_overlays.py <day> [deck.pptx]")
        sys.exit(1)
    day = int(sys.argv[1])
    folder = bd.DAY_FOLDERS[day]
    deck = HERE / "out" / (sys.argv[2] if len(sys.argv) > 2 else f"Day{day}.pptx")
    if not deck.exists():
        print(f"Deck not found: {deck}")
        sys.exit(1)

    # blueprint indexes: filename -> [slide_id], and slide_id -> identifying texts
    # (title + caption/prose). G09 "do-it" slides render prose-first with no title
    # textbox, so caption is needed to tell the where/do-it pair apart.
    slides = bd.parse_blueprint((bd.REPO / folder / "SLIDE_SOURCE.md").read_text(encoding="utf-8"))
    fname_to_slides = {}
    id_texts = {}
    for sl in slides:
        title = bd._clean(bd._first(sl["fields"], "Title")) or sl["label"]
        caption = bd._caption(sl["fields"])
        id_texts[sl["id"]] = {t for t in (title, caption) if t}
        gname = bd._resolved_guide_name(sl, day)
        if gname:
            fname_to_slides.setdefault(gname, []).append(sl["id"])

    disk = _disk_hash_index()
    prs = Presentation(str(deck))

    def _texts(slide):
        out = []
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                out.append(sh.text_frame.text.strip())
        return out

    def _id_matches(sid, old_texts):
        for idt in id_texts.get(sid, ()):
            if any(idt == ot or idt in ot for ot in old_texts):
                return True
        return False

    overlays = {}
    unmatched = []
    for i, slide in enumerate(prs.slides):
        boxes = [s for s in slide.shapes if _is_overlay(s)]
        if not boxes:
            continue
        geo = [{"left": int(s.left), "top": int(s.top),
                "width": int(s.width), "height": int(s.height)} for s in boxes]
        old_texts = _texts(slide)

        # 1) match by embedded screenshot hash; disambiguate shared shots by text
        sid = None
        for h in _slide_image_hashes(slide):
            fname = disk.get(h)
            if not fname:
                continue
            cands = fname_to_slides.get(fname, [])
            if len(cands) == 1:
                sid = cands[0]
            elif len(cands) > 1:
                sid = next((c for c in cands if _id_matches(c, old_texts)), cands[0])
            if sid:
                break

        # 2) fallback: unique identifying text (title or caption)
        if not sid:
            hits = [s2 for s2 in id_texts if _id_matches(s2, old_texts)]
            if len(hits) == 1:
                sid = hits[0]

        if sid:
            overlays.setdefault(sid, []).extend(geo)
        else:
            label = (_slide_title(slide) or (old_texts[0] if old_texts else ""))[:40]
            unmatched.append((i + 1, label, len(geo)))

    out = HERE / "out" / f"Day{day}.overlays.json"
    out.write_text(json.dumps(overlays, indent=2), encoding="utf-8")
    matched_boxes = sum(len(v) for v in overlays.values())
    print(f"Wrote {out}")
    print(f"  matched {matched_boxes} box(es) across {len(overlays)} slide(s)")
    if unmatched:
        print(f"  {len(unmatched)} slide(s) with boxes could NOT be matched:")
        for pg, t, n in unmatched:
            print(f"    old p{pg}: {n} box(es) — {t!r}")


if __name__ == "__main__":
    main()
