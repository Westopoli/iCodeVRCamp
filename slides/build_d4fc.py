"""Build D4FC.pptx — single reference card: all 3 FC examples on one slide.

D4 FC is creative: kids invent a 5th fighter. Three holes:
  FC-1 — fill the CUSTOM_CHARACTER stats dict
  FC-2 — register the character in main.gd
  FC-3 — add a "custom" branch in attack() (freeform)

Layout: 3 cols × 1 row (tall cells — FC-1 dict is long).

Usage:
    cd slides
    .venv/bin/python build_d4fc.py   -> out/D4FC.pptx
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import theme
import master

# ── data ──────────────────────────────────────────────────────────────────────

EXAMPLES = [
    {
        "label": "FC-1 — CUSTOM_CHARACTER dict",
        "mirror": "Chunks #1 + #2",
        "code": [
            "const CUSTOM_CHARACTER := {",
            '    "display_name": "MyFighter",',
            '    "sprite": "res://assets/.../tile_0004.png",',
            '    "tint": Color(1, 1, 1),',
            '    "walk_speed": 250.0,',
            '    "jump_impulse": 540.0,',
            '    "attack_type": "custom",',
            '    "attack_damage": 12,',
            '    "attack_cooldown": 0.6,',
            '    "attack_range": 0.0,',
            '    "projectile_speed": 0.0,',
            '    "projectile_gravity_scale": 0.0,',
            "}",
        ],
    },
    {
        "label": "FC-2 — register in main.gd",
        "mirror": "Chunk #4",
        "code": [
            "# In main.gd _ready():",
            'CHARACTERS["custom"] = CUSTOM_CHARACTER',
            "",
            "# Update char-select keys array:",
            'var keys = ["knight", "ninja",',
            '    "mage", "archer", "custom"]',
        ],
    },
    {
        "label": "FC-3 — custom attack() branch",
        "mirror": "Chunks #6 + #7",
        "code": [
            '# Invent anything — patterns available:',
            "opponent.take_damage(N)   # deal damage",
            "spawn_projectile()         # fire projectile",
            "hp += N                    # heal yourself",
            "melee_swing_timer = 0.15  # draw swing arc",
            "queue_redraw()",
            "",
            "# Multiple actions per attack are fine.",
            "# No fixed shape — this is your design.",
        ],
    },
]

# ── layout ────────────────────────────────────────────────────────────────────

COLS = 3
ROWS = 1
MARGIN_H = Inches(0.35)
MARGIN_TOP_BODY = Inches(0.18)
GAP_COL = Inches(0.22)
MARGIN_BOTTOM = Inches(0.18)

BODY_W = theme.SLIDE_WIDTH - MARGIN_H * 2
BODY_H = theme.SLIDE_HEIGHT - theme.HEADER_HEIGHT - MARGIN_TOP_BODY - MARGIN_BOTTOM

CELL_W = (BODY_W - GAP_COL * (COLS - 1)) / COLS
CELL_H = BODY_H  # single row — full body height

LABEL_H = Inches(0.40)
PAD = Inches(0.12)
CODE_SIZE = Pt(12)
LABEL_SIZE = Pt(12)

# ── tokeniser ─────────────────────────────────────────────────────────────────

_KEYWORDS = {
    "for", "in", "range", "var", "const", "func", "return",
    "if", "else", "elif", "and", "or", "not", "pass",
    "true", "false", "null",
}
_FLOW_KW = {"while", "if", "else", "elif", "match"}
_TYPES = {"bool", "int", "float", "String", "void", "Vector2i", "Vector2", "Array", "Node", "Color"}


def tokenize(line: str):
    if not line.strip():
        return [(" ", theme.CODE_TEXT)]
    result = []
    stripped = line.lstrip(" ")
    indent = len(line) - len(stripped)
    if indent:
        result.append((" " * indent, theme.CODE_TEXT))
    line = stripped
    i = 0
    while i < len(line):
        ch = line[i]
        if ch == "#":
            result.append((line[i:], theme.SYNTAX_COMMENT))
            break
        if ch in ('"', "'"):
            q = ch
            j = i + 1
            while j < len(line) and line[j] != q:
                j += 1
            j += 1
            result.append((line[i:j], theme.SYNTAX_STRING))
            i = j
            continue
        if ch.isdigit():
            j = i
            while j < len(line) and (line[j].isdigit() or line[j] == "."):
                j += 1
            result.append((line[i:j], theme.SYNTAX_NUMBER))
            i = j
            continue
        if ch.isalpha() or ch == "_":
            j = i
            while j < len(line) and (line[j].isalnum() or line[j] == "_"):
                j += 1
            word = line[i:j]
            if word in _TYPES:
                color = theme.SYNTAX_TYPE
            elif word in _FLOW_KW:
                color = theme.SYNTAX_KEYWORD_FLOW
            elif word in _KEYWORDS:
                color = theme.SYNTAX_KEYWORD
            elif j < len(line) and line[j] == "(":
                color = theme.SYNTAX_FUNCTION
            else:
                color = theme.CODE_TEXT
            result.append((word, color))
            i = j
            continue
        two = line[i:i+2]
        if two in ("->", ":=", "+=", "-=", "*=", "==", "!=", "<=", ">="):
            result.append((two, theme.SYNTAX_OPERATOR))
            i += 2
            continue
        result.append((ch, theme.CODE_TEXT))
        i += 1
    return result


# ── builder ───────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = theme.SLIDE_WIDTH
    prs.slide_height = theme.SLIDE_HEIGHT

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    master.apply_master(slide, day_number=4, page_number=None)

    title_box = slide.shapes.add_textbox(
        Inches(3.0), Inches(0.18), Inches(7.5), Inches(0.55),
    )
    tf = title_box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Final Challenge — Board Examples"
    run.font.name = theme.FONT_HEADING
    run.font.size = Pt(19)
    run.font.bold = True
    run.font.color.rgb = theme.BG_WHITE

    for idx, ex in enumerate(EXAMPLES):
        col = idx % COLS

        left = MARGIN_H + col * (CELL_W + GAP_COL)
        top = theme.HEADER_HEIGHT + MARGIN_TOP_BODY

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, CELL_W, CELL_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.CODE_BG
        bg.line.fill.background()
        bg.shadow.inherit = False

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, CELL_W, LABEL_H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme.ICODE_RED
        bar.line.fill.background()
        bar.shadow.inherit = False

        lbl = slide.shapes.add_textbox(left + PAD, top, CELL_W - PAD * 2, LABEL_H)
        tf = lbl.text_frame
        tf.margin_top = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        fc_num, _, label_rest = ex["label"].partition(" — ")
        run = p.add_run()
        run.text = fc_num + "  "
        run.font.name = theme.FONT_HEADING
        run.font.size = LABEL_SIZE
        run.font.bold = True
        run.font.color.rgb = theme.BG_WHITE
        run2 = p.add_run()
        run2.text = label_rest
        run2.font.name = theme.FONT_HEADING
        run2.font.size = LABEL_SIZE
        run2.font.bold = False
        run2.font.color.rgb = theme.BG_WHITE

        # Mirror label (small, below label bar)
        mirror_top = top + LABEL_H
        mirror_box = slide.shapes.add_textbox(
            left + PAD, mirror_top, CELL_W - PAD * 2, Inches(0.28),
        )
        tf = mirror_box.text_frame
        tf.margin_top = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"← {ex['mirror']}"
        run.font.name = theme.FONT_BODY
        run.font.size = Pt(10)
        run.font.italic = True
        run.font.color.rgb = theme.GREY_MID

        code_top = top + LABEL_H + Inches(0.28) + PAD
        code_h = CELL_H - LABEL_H - Inches(0.28) - PAD * 2
        code_box = slide.shapes.add_textbox(
            left + PAD, code_top, CELL_W - PAD * 2, code_h
        )
        tf = code_box.text_frame
        tf.margin_top = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.word_wrap = False

        for li, line in enumerate(ex["code"]):
            para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            for text, color in tokenize(line):
                run = para.add_run()
                run.text = text
                run.font.name = theme.FONT_MONO
                run.font.size = CODE_SIZE
                run.font.color.rgb = color

    out = Path(__file__).parent / "out" / "D4FC.pptx"
    prs.save(str(out))
    print(f"Wrote {out} — 1 slide")


if __name__ == "__main__":
    build()
