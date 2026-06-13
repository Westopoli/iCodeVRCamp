"""Build D3FC.pptx — single reference card: all 9 FC board examples on one slide.

Usage:
    cd slides
    .venv/bin/python build_d3fc.py   -> out/D3FC.pptx
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import theme
import master

# ── data ──────────────────────────────────────────────────────────────────────

EXAMPLES = [
    {
        "label": "FC-1 — state vars",
        "mirror": "Chunk #1",
        "code": [
            "var spawn_timer: float = 0.0",
            "var difficulty: int = 1",
            "var spawn_interval: float = SPAWN_INTERVAL_START",
            "var spawn_queue: Array = []",
            "var clear_streak: int = 0",
        ],
    },
    {
        "label": "FC-2a — .append()",
        "mirror": "Chunk #2a",
        "code": [
            "spawn_queue.append(t)",
        ],
    },
    {
        "label": "FC-2b — .pop_front() + reward",
        "mirror": "Chunk #2b",
        "code": [
            "var t: String = spawn_queue.pop_front()",
            "main.spawn_enemy(random_edge(), t)",
            "main.coins += STREAK_BONUS",
        ],
    },
    {
        "label": "FC-3 — per-frame buff sweep",
        "mirror": "Chunk #3",
        "code": [
            "for e in main.enemies:",
            "    endless_buff(e, delta)",
            "for t in main.towers:",
            "    buff_tower(t, delta)",
        ],
    },
    {
        "label": "FC-4 — buff_all(list, delta)",
        "mirror": "Chunk #4",
        "code": [
            "for e in enemy_list:",
            "    endless_buff(e, delta)",
        ],
    },
    {
        "label": "FC-5a — get_fastest_enemy()",
        "mirror": "Chunk #5a",
        "code": [
            "for e in main.enemies:",
            "    if e.speed > best_speed:",
            "        fastest = e",
            "        best_speed = e.speed",
        ],
    },
    {
        "label": "FC-5b — get_wounded_enemies()",
        "mirror": "Chunk #5b",
        "code": [
            "var result: Array = []",
            "for e in main.enemies:",
            "    if e.hp <= WOUNDED_HP_THRESHOLD:",
            "        result.append(e)",
            "return result",
        ],
    },
    {
        "label": "FC-6 — escalate() match branch",
        "mirror": "Chunk #6",
        "code": [
            '"medium":',
            "    var t: String = pick_type_for_band(band)",
            "    queue_spawn(t)",
            "    queue_spawn(t)",
        ],
    },
    {
        "label": "FC-7 — check_for_screen_clear()",
        "mirror": "Chunk #7",
        "code": [
            "if main.enemies.size() == 0 and spawn_queue.size() == 0:",
            "    clear_streak += 1",
            "    difficulty += 1",
            "    spawn_interval *= SPAWN_INTERVAL_SHRINK",
            "    escalate()",
        ],
    },
]

# ── layout ────────────────────────────────────────────────────────────────────

COLS = 3
ROWS = 3
MARGIN_H = Inches(0.35)
MARGIN_TOP_BODY = Inches(0.18)
GAP_COL = Inches(0.18)
GAP_ROW = Inches(0.18)
MARGIN_BOTTOM = Inches(0.18)

BODY_W = theme.SLIDE_WIDTH - MARGIN_H * 2
BODY_H = theme.SLIDE_HEIGHT - theme.HEADER_HEIGHT - MARGIN_TOP_BODY - MARGIN_BOTTOM

CELL_W = (BODY_W - GAP_COL * (COLS - 1)) / COLS
CELL_H = (BODY_H - GAP_ROW * (ROWS - 1)) / ROWS

LABEL_H = Inches(0.34)
PAD = Inches(0.09)
CODE_SIZE = Pt(10.5)
LABEL_SIZE = Pt(11)

# ── tokeniser (shared with d2fc) ──────────────────────────────────────────────

_KEYWORDS = {
    "for", "in", "range", "var", "func", "return",
    "if", "else", "elif", "and", "or", "not", "pass",
    "true", "false", "null",
}
_FLOW_KW = {"while", "if", "else", "elif", "match"}
_TYPES = {"bool", "int", "float", "String", "void", "Vector2i", "Vector2", "Array", "Node"}


def tokenize(line: str):
    result = []
    stripped = line.lstrip(" ")
    indent = len(line) - len(stripped)
    if indent:
        result.append((" " * indent, theme.CODE_TEXT))
    line = stripped
    i = 0
    while i < len(line):
        ch = line[i]
        if ch == "#":
            result.append((line[i:], theme.SYNTAX_COMMENT))
            break
        if ch in ('"', "'"):
            q = ch
            j = i + 1
            while j < len(line) and line[j] != q:
                j += 1
            j += 1
            result.append((line[i:j], theme.SYNTAX_STRING))
            i = j
            continue
        if ch.isdigit():
            j = i
            while j < len(line) and (line[j].isdigit() or line[j] == "."):
                j += 1
            result.append((line[i:j], theme.SYNTAX_NUMBER))
            i = j
            continue
        if ch.isalpha() or ch == "_":
            j = i
            while j < len(line) and (line[j].isalnum() or line[j] == "_"):
                j += 1
            word = line[i:j]
            if word in _TYPES:
                color = theme.SYNTAX_TYPE
            elif word in _FLOW_KW:
                color = theme.SYNTAX_KEYWORD_FLOW
            elif word in _KEYWORDS:
                color = theme.SYNTAX_KEYWORD
            elif j < len(line) and line[j] == "(":
                color = theme.SYNTAX_FUNCTION
            else:
                color = theme.CODE_TEXT
            result.append((word, color))
            i = j
            continue
        two = line[i:i+2]
        if two in ("->", ":=", "+=", "-=", "*=", "==", "!=", "<=", ">="):
            result.append((two, theme.SYNTAX_OPERATOR))
            i += 2
            continue
        result.append((ch, theme.CODE_TEXT))
        i += 1
    return result


# ── builder ───────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = theme.SLIDE_WIDTH
    prs.slide_height = theme.SLIDE_HEIGHT

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    master.apply_master(slide, day_number=3, page_number=None)

    title_box = slide.shapes.add_textbox(
        Inches(3.0), Inches(0.18), Inches(7.5), Inches(0.55),
    )
    tf = title_box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Final Challenge — Board Examples"
    run.font.name = theme.FONT_HEADING
    run.font.size = Pt(19)
    run.font.bold = True
    run.font.color.rgb = theme.BG_WHITE

    for idx, ex in enumerate(EXAMPLES):
        col = idx % COLS
        row = idx // COLS

        left = MARGIN_H + col * (CELL_W + GAP_COL)
        top = theme.HEADER_HEIGHT + MARGIN_TOP_BODY + row * (CELL_H + GAP_ROW)

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, CELL_W, CELL_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.CODE_BG
        bg.line.fill.background()
        bg.shadow.inherit = False

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, CELL_W, LABEL_H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme.ICODE_RED
        bar.line.fill.background()
        bar.shadow.inherit = False

        lbl = slide.shapes.add_textbox(left + PAD, top, CELL_W - PAD * 2, LABEL_H)
        tf = lbl.text_frame
        tf.margin_top = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        fc_num, _, label_rest = ex["label"].partition(" — ")
        run = p.add_run()
        run.text = fc_num + "  "
        run.font.name = theme.FONT_HEADING
        run.font.size = LABEL_SIZE
        run.font.bold = True
        run.font.color.rgb = theme.BG_WHITE
        run2 = p.add_run()
        run2.text = label_rest
        run2.font.name = theme.FONT_HEADING
        run2.font.size = LABEL_SIZE
        run2.font.bold = False
        run2.font.color.rgb = theme.BG_WHITE

        code_top = top + LABEL_H + PAD
        code_h = CELL_H - LABEL_H - PAD * 2
        code_box = slide.shapes.add_textbox(
            left + PAD, code_top, CELL_W - PAD * 2, code_h
        )
        tf = code_box.text_frame
        tf.margin_top = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.word_wrap = False

        for li, line in enumerate(ex["code"]):
            para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            for text, color in tokenize(line):
                run = para.add_run()
                run.text = text
                run.font.name = theme.FONT_MONO
                run.font.size = CODE_SIZE
                run.font.color.rgb = color

    out = Path(__file__).parent / "out" / "D3FC.pptx"
    prs.save(str(out))
    print(f"Wrote {out} — 1 slide")


if __name__ == "__main__":
    build()
