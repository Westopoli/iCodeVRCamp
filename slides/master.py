"""Master-frame helpers used by every slide layout.

Each slide gets:
- A solid black header bar at the top (red/black/grey minimalist brand)
- The iCode logo top-left (on the bar)
- A red "Day N" label top-right (no colored tab — text is the only per-day mark)
- A page number bottom-right

`apply_master(slide, day_number, page_number)` is called as the FIRST step of
every layout function in `templates.py`.
"""

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE

import theme


def _add_header_bar(slide):
    """Add the full-width solid black bar across the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        theme.SLIDE_WIDTH, theme.HEADER_HEIGHT,
    )
    shape.line.fill.background()  # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.HEADER_BG
    shape.shadow.inherit = False
    return shape


def _add_logo(slide):
    """Drop the iCode logo top-left on the header strip + the iCode wordmark
    + tagline next to it (matches the brand sample decks)."""
    logo_height = Inches(0.65)
    logo_top = Inches(0.15)
    logo_left = Inches(0.3)

    if theme.LOGO_RED_PATH.exists():
        slide.shapes.add_picture(
            str(theme.LOGO_RED_PATH),
            logo_left, logo_top,
            height=logo_height,
        )
        # Match the logo's natural aspect (~1.3:1) so the text-box starts at the right edge
        wordmark_left = logo_left + Inches(0.85)
    else:
        wordmark_left = logo_left

    # "iCode" wordmark — large bold white
    wordmark_box = slide.shapes.add_textbox(
        wordmark_left, Inches(0.05),
        Inches(2.0), Inches(0.55),
    )
    tf = wordmark_box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = 1  # left
    run = p.add_run()
    run.text = "iCode"
    run.font.name = theme.FONT_HEADING
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = theme.BG_WHITE

    # Tagline — small white below the wordmark
    tagline_box = slide.shapes.add_textbox(
        wordmark_left, Inches(0.6),
        Inches(2.5), Inches(0.25),
    )
    tf = tagline_box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.add_run()
    run.text = "EMPOWERING FUTURE INNOVATORS"
    run.font.name = theme.FONT_BODY
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = theme.BG_WHITE


def _add_day_label(slide, day_number):
    """Add a minimalist red 'DAY N' label top-right on the black bar. No tab
    shape — the red text is the only per-day differentiator."""
    if day_number is None:
        return
    width = Inches(1.6)
    height = Inches(0.45)
    left = theme.SLIDE_WIDTH - width - Inches(0.3)
    top = Inches(0.25)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = 2  # right
    run = p.add_run()
    run.text = f"DAY {day_number}"
    run.font.name = theme.FONT_HEADING
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = theme.ICODE_RED


def _add_page_number(slide, page_number):
    """Drop a page number bottom-right."""
    if page_number is None:
        return
    width = Inches(0.8)
    height = Inches(0.3)
    left = theme.SLIDE_WIDTH - width - Inches(0.2)
    top = theme.SLIDE_HEIGHT - height - Inches(0.15)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = 2
    run = p.add_run()
    run.text = str(page_number)
    run.font.name = theme.FONT_BODY
    run.font.size = theme.SIZE_PAGE_NUMBER
    run.font.color.rgb = theme.TEXT_MUTED


def apply_master(slide, day_number=None, page_number=None):
    """Apply the master frame (header strip + logo + day tab + page number) to a slide.

    Call this as the FIRST step of every layout function in templates.py.
    """
    _add_header_bar(slide)
    _add_logo(slide)
    _add_day_label(slide, day_number)
    _add_page_number(slide, page_number)
