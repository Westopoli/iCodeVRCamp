"""Build D4FCbeta_v2.pptx — 2 slides: the 6 steps to add a 5th character.

Base game has 4 characters (knight, ninja, mage, archer). Students add the 5th.
Follows build_fc_beta.py card layout exactly (3 cols x 1 row per slide).
Output: slides/out/D4FCbeta_v2.pptx

Usage:
    cd slides
    python build_d4fcbeta.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ---- Brand colors (mirrors theme.py) ----
BLACK    = RGBColor(0x11, 0x11, 0x11)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
NOTE_GRY = RGBColor(0x8A, 0x8A, 0x8A)
BOX_BDR  = RGBColor(0xCC, 0xCC, 0xCC)
CODE_BG  = RGBColor(0x2D, 0x2D, 0x2D)
CODE_LT  = RGBColor(0xE6, 0xE6, 0xE6)

FONT_UI   = "Poppins"
FONT_MONO = "Consolas"

SW = Inches(13.333)
SH = Inches(7.5)


# =====================================================================
#  Slide data
# =====================================================================

SLIDE_1_BOXES = [
    (
        "Step 1 — Add entry to CHARACTERS  (main.gd)",
        "melee needs attack_range > 0  |  projectile needs projectile_speed > 0",
        [
            '"yourchar": {',
            '    "display_name": "YourChar",',
            '    "sprite": "res://assets/.../tile_0003.png",',
            '    "tint": Color(1.0, 0.5, 0.0),',
            '    "walk_speed": 250.0,',
            '    "jump_impulse": 500.0,',
            '    "attack_type": "melee",',
            '    "attack_damage": 20,',
            '    "attack_cooldown": 0.6,',
            '    "attack_range": 70.0,',
            '    "projectile_speed": 0.0,',
            '    "projectile_gravity_scale": 0.0,',
            '},',
        ],
    ),
    (
        "Step 2 — Add name to keys list  (main.gd)",
        "_unhandled_input()  •  name must match Step 1 key exactly",
        [
            "var keys = [",
            '    "knight", "ninja",',
            '    "mage", "archer",',
            '    "yourchar"',
            "]",
        ],
    ),
    (
        "Step 3 — Expand key capture range  (main.gd)",
        "_unhandled_input()  •  change the upper bound from 4 to 5",
        [
            "# Before (captures keys 1-4):",
            "event.keycode <= KEY_4",
            "",
            "# After (captures keys 1-5):",
            "event.keycode <= KEY_5",
        ],
    ),
]

SLIDE_2_BOXES = [
    (
        "Step 4 — Update selection guard  (main.gd)",
        "_unhandled_input()  •  change in BOTH char_select_p1 and char_select_p2 branches",
        [
            "# Before (both char_select branches):",
            "if key_num >= 1 and key_num <= 4:",
            "    p1_choice = keys[key_num - 1]",
            "",
            "# After:",
            "if key_num >= 1 and key_num <= 5:",
            "    p1_choice = keys[key_num - 1]",
        ],
    ),
    (
        "Step 5 — Update menu display text  (main.gd)",
        'set_screen()  •  update BOTH "char_select_p1" and "char_select_p2" text strings',
        [
            "# Find title_label.text in set_screen().",
            "# Add your character at the end:",
            "",
            '"P1 — pick your fighter:\\n"',
            '"1=Knight  2=Ninja  3=Mage\\n"',
            '"4=Archer  5=YourChar\\n"',
            '"(Space confirms)"',
        ],
    ),
    (
        "Step 6 — Add attack branch  (player.gd)",
        'Only needed if attack_type is "custom" — skip for melee or projectile',
        [
            '# In attack()  match attack_type:',
            '"custom":',
            '    var opp = get_opponent()',
            '    if opp == null or opp.is_dead():',
            '        return',
            '    opp.take_damage(attack_damage)',
            '    # hp += 10   # heal instead',
            '    # spawn_projectile()   # fire shot',
        ],
    ),
]

SLIDES = [
    {
        "title": "Day 4  --  Adding a 5th Character",
        "subtitle": "Steps 1-3 of 6  --  main.gd",
        "boxes": SLIDE_1_BOXES,
        "cols": 3,
        "rows": 1,
    },
    {
        "title": "Day 4  --  Adding a 5th Character",
        "subtitle": "Steps 4-6 of 6  --  main.gd  /  player.gd",
        "boxes": SLIDE_2_BOXES,
        "cols": 3,
        "rows": 1,
    },
]


# =====================================================================
#  Helpers  (mirrors build_fc_beta.py)
# =====================================================================

def set_run(run, text, font_name, size_pt, color, bold=False, italic=False):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def shape_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def shape_no_line(shape):
    shape.line.fill.background()


def shape_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


# =====================================================================
#  Card builder  (mirrors build_fc_beta.py add_card)
# =====================================================================

def add_card(slide, x, y, w, h, label, note, code_lines, code_pt=7.0):
    PAD          = Inches(0.10)
    BAR_H        = Inches(0.22)
    CODE_TOP_PAD = Inches(0.03)
    CODE_BOT_PAD = Inches(0.04)
    NOTE_H       = Inches(0.20)

    card = slide.shapes.add_shape(1, x, y, w, h)
    shape_fill(card, WHITE)
    shape_line(card, BOX_BDR, 0.75)

    bar = slide.shapes.add_shape(1, x, y, w, BAR_H)
    shape_fill(bar, BLACK)
    shape_no_line(bar)

    lbl_tb = slide.shapes.add_textbox(
        x + PAD, y + Inches(0.035), w - 2 * PAD, BAR_H - Inches(0.04)
    )
    tf = lbl_tb.text_frame
    tf.word_wrap = False
    run = tf.paragraphs[0].add_run()
    set_run(run, label, FONT_UI, 6.8, WHITE, bold=True)

    code_y = y + BAR_H + CODE_TOP_PAD
    line_h_in = (code_pt * 1.35) / 72.0
    code_h = Inches(line_h_in * len(code_lines) + 0.05)
    max_code_h = h - BAR_H - CODE_TOP_PAD - CODE_BOT_PAD - NOTE_H - Inches(0.05)
    code_h = min(code_h, max_code_h)

    code_bg = slide.shapes.add_shape(
        1, x + PAD - Inches(0.06), code_y, w - 2 * PAD + Inches(0.12), code_h
    )
    shape_fill(code_bg, CODE_BG)
    shape_no_line(code_bg)

    code_tb = slide.shapes.add_textbox(
        x + PAD, code_y + Inches(0.025), w - 2 * PAD, code_h
    )
    ctf = code_tb.text_frame
    ctf.word_wrap = False
    first = True
    for line in code_lines:
        p = ctf.paragraphs[0] if first else ctf.add_paragraph()
        first = False
        p.space_before = Pt(0)
        run = p.add_run()
        set_run(run, line, FONT_MONO, code_pt, CODE_LT)

    note_y = code_y + code_h + CODE_BOT_PAD
    note_tb = slide.shapes.add_textbox(x + PAD, note_y, w - 2 * PAD, NOTE_H)
    ntf = note_tb.text_frame
    ntf.word_wrap = True
    run = ntf.paragraphs[0].add_run()
    set_run(run, note, FONT_UI, 5.5, NOTE_GRY, italic=True)


# =====================================================================
#  Slide builder  (mirrors build_fc_beta.py build_slide)
# =====================================================================

def build_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    shape_fill(bg, WHITE)
    shape_no_line(bg)

    TITLE_H = Inches(0.50)
    title_bar = slide.shapes.add_shape(1, 0, 0, SW, TITLE_H)
    shape_fill(title_bar, BLACK)
    shape_no_line(title_bar)

    tb = slide.shapes.add_textbox(Inches(0.28), Inches(0.08), Inches(12), Inches(0.38))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run()
    set_run(r1, slide_data["title"], FONT_UI, 15, WHITE, bold=True)
    r2 = p.add_run()
    set_run(r2, "   " + slide_data["subtitle"], FONT_UI, 10, NOTE_GRY)

    COLS = slide_data["cols"]
    ROWS = slide_data["rows"]
    boxes = slide_data["boxes"]

    ML = Inches(0.17)
    MT = TITLE_H + Inches(0.11)
    MB = Inches(0.11)
    GH = Inches(0.08)
    GV = Inches(0.08)

    total_w = SW - 2 * ML
    total_h = SH - MT - MB
    card_w = (total_w - GH * (COLS - 1)) / COLS
    card_h = (total_h - GV * (ROWS - 1)) / ROWS

    max_lines = max(len(b[2]) for b in boxes)
    avail_h_in = (card_h / 914400) - 0.22 - 0.20 - 0.10
    code_pt = min(8.0, avail_h_in * 72.0 / (max_lines * 1.35))
    code_pt = max(5.5, code_pt)

    for idx, (label, note, code_lines) in enumerate(boxes):
        row = idx // COLS
        col = idx % COLS
        x = ML + col * (card_w + GH)
        y = MT + row * (card_h + GV)
        add_card(
            slide,
            int(x), int(y), int(card_w), int(card_h),
            label, note, code_lines,
            code_pt=code_pt,
        )


# =====================================================================
#  Main
# =====================================================================

def main():
    out_dir = Path(__file__).parent / "out"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "D4FCbeta_v2.pptx"

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    for slide_data in SLIDES:
        build_slide(prs, slide_data)

    prs.save(str(out_path))
    print("Saved: " + str(out_path))
    print("  Slide 1: Steps 1-3  (4 base chars -> student adds 5th)")
    print("  Slide 2: Steps 4-6")


if __name__ == "__main__":
    main()
