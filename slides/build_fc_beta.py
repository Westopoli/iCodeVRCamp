"""
Build FC beta syntax-reference slides for D2, D3, D4.
Output: slides/out/FC_beta.pptx  (one slide per day)

Each slide is a grid of cards -- one card per FC hole.
Each card shows: FC label (white on black bar), literal code (dark on dark bg),
and a small grey note linking to the morning chunk.

Usage:
    cd slides
    python build_fc_beta.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ---- Brand colors (from theme.py) ----
BLACK     = RGBColor(0x11, 0x11, 0x11)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
NOTE_GRY  = RGBColor(0x8A, 0x8A, 0x8A)
BOX_BDR   = RGBColor(0xCC, 0xCC, 0xCC)
CODE_BG   = RGBColor(0x2D, 0x2D, 0x2D)
CODE_LT   = RGBColor(0xE6, 0xE6, 0xE6)

# ---- Fonts ----
FONT_UI   = "Poppins"
FONT_MONO = "Consolas"

# ---- Slide geometry ----
SW = Inches(13.333)
SH = Inches(7.5)


# =====================================================================
#  FC DATA -- (label, note, code_lines)
#  label : shown white-on-black in the card header bar
#  note  : small grey text at bottom of card, links to morning chunk
#  code_lines : literal GDScript, actual names from the game file
# =====================================================================

D2_BOXES = [
    (
        "FC-1  spawn_personality_ghosts()",
        "mirrors Chunk #1",
        [
            "for i in range(PERSONALITY_COUNT):",
            "    spawn_one_personality(i)",
        ],
    ),
    (
        "FC-2  step_all_personality_ghosts()",
        "mirrors Chunk #2",
        [
            "for ghost in ghosts:",
            "    step_personality(ghost)",
        ],
    ),
    (
        "FC-3  count_ghosts_of(personality) -> int",
        "mirrors Chunk #3  --  get_meta() reads the tag",
        [
            "var count := 0",
            "for ghost in ghosts:",
            '    if ghost.get_meta("personality") == personality:',
            "        count += 1",
            "return count",
        ],
    ),
    (
        "FC-4  reset_personality_ghosts()",
        "mirrors Chunk #4  --  i tracks the slot index",
        [
            "var i := 0",
            "for ghost in ghosts:",
            "    respawn_personality_ghost(ghost, i)",
            "    i += 1",
        ],
    ),
    (
        "FC-5  target_for(ghost) -> Vector2i",
        "mirrors Chunk #5  --  get_meta + return inside if",
        [
            'var p = ghost.get_meta("personality")',
            "if p == BLINKY:  return blinky_target(ghost)",
            "if p == PINKY:   return pinky_target(ghost)",
            "if p == INKY:    return inky_target(ghost)",
            "return clyde_target(ghost)",
        ],
    ),
    (
        "FC-6  is_clyde_close(ghost) -> bool",
        "mirrors Chunk #6  --  return a comparison directly",
        [
            "var d := distance_to_player(ghost)",
            "return d < 8.0",
        ],
    ),
]

D3_BOXES = [
    (
        "FC-1  state variables",
        "mirrors Chunk #1",
        [
            "var spawn_timer: float = 0.0",
            "var difficulty: int = 1",
            "var spawn_interval: float = SPAWN_INTERVAL_START",
            "var spawn_queue: Array = []",
            "var clear_streak: int = 0",
        ],
    ),
    (
        "FC-2a  queue_spawn(t)",
        "mirrors Chunk #2a  --  .append() adds to back of list",
        [
            "spawn_queue.append(t)",
        ],
    ),
    (
        "FC-2b  take_next_spawn()",
        "mirrors Chunk #2b  --  .pop_front() removes from front",
        [
            "var t: String = spawn_queue.pop_front()",
            "main.spawn_enemy(random_edge(), t)",
            "main.coins += STREAK_BONUS",
        ],
    ),
    (
        "FC-3  per-frame buff sweep  (inside endless_tick)",
        "mirrors Chunk #3  --  two for-loops back to back",
        [
            "for e in main.enemies:",
            "    endless_buff(e, delta)",
            "for t in main.towers:",
            "    buff_tower(t, delta)",
        ],
    ),
    (
        "FC-4  buff_all(enemy_list, delta)",
        "mirrors Chunk #4  --  list arrives as a parameter",
        [
            "for e in enemy_list:",
            "    endless_buff(e, delta)",
        ],
    ),
    (
        "FC-5a  get_fastest_enemy() -> Node",
        "mirrors Chunk #5a  --  track best as you loop",
        [
            "for e in main.enemies:",
            "    if e.speed > best_speed:",
            "        fastest = e",
            "        best_speed = e.speed",
        ],
    ),
    (
        "FC-5b  get_wounded_enemies() -> Array",
        "mirrors Chunk #5b  --  build + return a new list",
        [
            "var result: Array = []",
            "for e in main.enemies:",
            "    if e.hp <= WOUNDED_HP_THRESHOLD:",
            "        result.append(e)",
            "return result",
        ],
    ),
    (
        "FC-6  escalate()  match branches",
        "mirrors Chunk #6  --  repeat queue_spawn per difficulty",
        [
            "var t = pick_type_for_band(band)",
            'queue_spawn(t)        # easy:   x1',
            'queue_spawn(t)        # medium: x2',
            '# hard: x3   insane: x4',
        ],
    ),
    (
        "FC-7  check_for_screen_clear()",
        "mirrors Chunk #7  --  .size()==0 on both lists",
        [
            "if main.enemies.size() == 0 \\",
            "        and spawn_queue.size() == 0:",
            "    clear_streak += 1",
            "    difficulty += 1",
            "    spawn_interval *= SPAWN_INTERVAL_SHRINK",
            "    main.base_hp += BASE_HP_REGEN_PER_CLEAR",
            "    escalate()",
        ],
    ),
]

D4_BOXES = [
    (
        "FC-1  Fill your character's stats",
        'Change any values -- keep  attack_type: "custom"',
        [
            "const CUSTOM_CHARACTER := {",
            '    "display_name": "YourName",',
            '    "walk_speed":   250.0,',
            '    "jump_impulse": 540.0,',
            '    "attack_damage": 12,',
            '    "attack_type": "custom",',
            "}",
        ],
    ),
    (
        "FC-2  Register in main.gd",
        "In main.gd _ready(), after CHARACTERS is defined",
        [
            'CHARACTERS["custom"] = CUSTOM_CHARACTER',
        ],
    ),
    (
        'FC-3  Custom attack  (player.gd)',
        'Add a "custom": branch in the attack() match statement',
        [
            '"custom":',
            "    opponent.take_damage(attack_damage)",
            "    # examples:",
            "    # hp += 10              -> heal yourself",
            "    # spawn_projectile()    -> fire a shot",
            "    # take_damage() twice   -> double hit",
        ],
    ),
]

DAYS = [
    {
        "title": "Day 2  --  Final Challenge Syntax",
        "subtitle": "ghost_personalities.gd",
        "boxes": D2_BOXES,
        "cols": 3,
        "rows": 2,
    },
    {
        "title": "Day 3  --  Final Challenge Syntax",
        "subtitle": "endless_mode.gd",
        "boxes": D3_BOXES,
        "cols": 3,
        "rows": 3,
    },
    {
        "title": "Day 4  --  Final Challenge Syntax",
        "subtitle": "final_challenge.gd  /  player.gd  /  main.gd",
        "boxes": D4_BOXES,
        "cols": 3,
        "rows": 1,
    },
]


# =====================================================================
#  Helpers
# =====================================================================

def set_run(run, text, font_name, size_pt, color, bold=False, italic=False):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def shape_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def shape_no_line(shape):
    shape.line.fill.background()


def shape_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


# =====================================================================
#  Card builder
# =====================================================================

def add_card(slide, x, y, w, h, label, note, code_lines, code_pt=7.0):
    """Draw one FC card at position (x, y) with size (w, h) -- all EMU."""
    PAD = Inches(0.10)
    BAR_H = Inches(0.22)
    CODE_TOP_PAD = Inches(0.03)
    CODE_BOT_PAD = Inches(0.04)
    NOTE_H = Inches(0.20)

    # outer card: white fill + grey border
    card = slide.shapes.add_shape(1, x, y, w, h)
    shape_fill(card, WHITE)
    shape_line(card, BOX_BDR, 0.75)

    # label bar: black stripe at top
    bar = slide.shapes.add_shape(1, x, y, w, BAR_H)
    shape_fill(bar, BLACK)
    shape_no_line(bar)

    # label text in bar
    lbl_tb = slide.shapes.add_textbox(
        x + PAD, y + Inches(0.035), w - 2 * PAD, BAR_H - Inches(0.04)
    )
    tf = lbl_tb.text_frame
    tf.word_wrap = False
    run = tf.paragraphs[0].add_run()
    set_run(run, label, FONT_UI, 6.8, WHITE, bold=True)

    # dark code area background
    code_y = y + BAR_H + CODE_TOP_PAD
    line_h_in = (code_pt * 1.35) / 72.0
    code_h = Inches(line_h_in * len(code_lines) + 0.05)
    max_code_h = h - BAR_H - CODE_TOP_PAD - CODE_BOT_PAD - NOTE_H - Inches(0.05)
    code_h = min(code_h, max_code_h)

    code_bg_shape = slide.shapes.add_shape(
        1, x + PAD - Inches(0.06), code_y, w - 2 * PAD + Inches(0.12), code_h
    )
    shape_fill(code_bg_shape, CODE_BG)
    shape_no_line(code_bg_shape)

    # code text
    code_tb = slide.shapes.add_textbox(
        x + PAD, code_y + Inches(0.025), w - 2 * PAD, code_h
    )
    ctf = code_tb.text_frame
    ctf.word_wrap = False
    first = True
    for line in code_lines:
        if first:
            p = ctf.paragraphs[0]
            first = False
        else:
            p = ctf.add_paragraph()
        p.space_before = Pt(0)
        run = p.add_run()
        set_run(run, line, FONT_MONO, code_pt, CODE_LT)

    # note text
    note_y = code_y + code_h + CODE_BOT_PAD
    note_tb = slide.shapes.add_textbox(
        x + PAD, note_y, w - 2 * PAD, NOTE_H
    )
    ntf = note_tb.text_frame
    ntf.word_wrap = True
    run = ntf.paragraphs[0].add_run()
    set_run(run, note, FONT_UI, 5.5, NOTE_GRY, italic=True)


# =====================================================================
#  Slide builder
# =====================================================================

def build_slide(prs, day):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # white background
    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    shape_fill(bg, WHITE)
    shape_no_line(bg)

    # title bar
    TITLE_H = Inches(0.50)
    title_bar = slide.shapes.add_shape(1, 0, 0, SW, TITLE_H)
    shape_fill(title_bar, BLACK)
    shape_no_line(title_bar)

    # title text
    tb = slide.shapes.add_textbox(Inches(0.28), Inches(0.08), Inches(12), Inches(0.38))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run()
    set_run(r1, day["title"], FONT_UI, 15, WHITE, bold=True)
    r2 = p.add_run()
    set_run(r2, "   " + day["subtitle"], FONT_UI, 10, NOTE_GRY)

    # grid layout
    COLS = day["cols"]
    ROWS = day["rows"]
    boxes = day["boxes"]

    ML = Inches(0.17)
    MT = TITLE_H + Inches(0.11)
    MB = Inches(0.11)
    GH = Inches(0.08)
    GV = Inches(0.08)

    total_w = SW - 2 * ML
    total_h = SH - MT - MB
    card_w = (total_w - GH * (COLS - 1)) / COLS
    card_h = (total_h - GV * (ROWS - 1)) / ROWS

    # auto-size code font to fit max lines in card height
    max_lines = max(len(b[2]) for b in boxes)
    avail_h_in = (card_h / 914400) - 0.22 - 0.20 - 0.10
    code_pt = min(8.0, avail_h_in * 72.0 / (max_lines * 1.35))
    code_pt = max(5.5, code_pt)

    for idx, (label, note, code_lines) in enumerate(boxes):
        row = idx // COLS
        col = idx % COLS
        x = ML + col * (card_w + GH)
        y = MT + row * (card_h + GV)
        add_card(
            slide,
            int(x), int(y), int(card_w), int(card_h),
            label, note, code_lines,
            code_pt=code_pt,
        )


# =====================================================================
#  Main
# =====================================================================

def main():
    out_dir = Path(__file__).parent / "out"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "FC_beta.pptx"

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    for day in DAYS:
        build_slide(prs, day)

    prs.save(str(out_path))
    print("Saved: " + str(out_path))


if __name__ == "__main__":
    main()
