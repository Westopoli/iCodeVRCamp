"""Generate standalone FC reference slide for each day (D1FC–D4FC).

Layout matches Day2v2 FC slide: 10"×5.625", 3-column box grid.
Each box: label strip | syntax strip | comment scaffold (dark bg, green text).

Usage (from slides/):
    .venv/bin/python gen_fc_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── Colors ───────────────────────────────────────────────────────────────────
BAR_BLACK     = RGBColor(0x11, 0x11, 0x11)
ICODE_RED     = RGBColor(0xE5, 0x3A, 0x2C)
GREY_DARK     = RGBColor(0x2B, 0x2B, 0x2B)
GREY_MID      = RGBColor(0x8A, 0x8A, 0x8A)
GREY_LIGHT    = RGBColor(0xF2, 0xF2, 0xF2)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_BLACK    = RGBColor(0x1A, 0x1A, 0x1A)
CODE_BG       = RGBColor(0x2D, 0x2D, 0x2D)
CODE_TEXT     = RGBColor(0xE6, 0xE6, 0xE6)
SYNTAX_COMMENT = RGBColor(0x6A, 0x99, 0x55)  # muted green (comment color)
SYNTAX_KEY    = RGBColor(0x56, 0x9C, 0xD6)   # blue (keyword)
LABEL_BG      = RGBColor(0x1A, 0x1A, 0x1A)   # near-black label strip
SYNTAX_BG     = RGBColor(0x3A, 0x3A, 0x3A)   # slightly lighter syntax strip

FONT_HEAD = "Poppins"
FONT_MONO = "Consolas"

# ── Slide geometry ────────────────────────────────────────────────────────────
SW = 9144000   # 10" × 914400
SH = 5143500   # 5.625"
HEADER_H = 460000
TITLE_H  = 230000
CONTENT_TOP = HEADER_H + TITLE_H
CONTENT_BOT = SH - 80000
CONTENT_H = CONTENT_BOT - CONTENT_TOP

LOGO_PATH = Path(__file__).parent / "assets" / "logos" / "icode_logo_red.png"

# ── FC content data ───────────────────────────────────────────────────────────
FC_DATA = {
    1: {
        "title": "Final Challenge — 3 holes",
        "day": "DAY 1",
        "grid": "3x1",
        "fcs": [
            {
                "label": "FC-1  Comment out AI",
                "syntax": "comment",
                "code": [
                    "# Find the 5 AI lines in _process()",
                    "# (they move right_paddle toward ball.position)",
                    "# Select each line and press Ctrl+K",
                    "# to comment it out",
                    "# Paddle should now sit still",
                ],
            },
            {
                "label": "FC-2  I/K key movement",
                "syntax": "if, input_key, minus_eq, plus_eq",
                "code": [
                    "# if I key pressed:",
                    "#     position.y -= paddle_speed",
                    "# if K key pressed:",
                    "#     position.y += paddle_speed",
                ],
            },
            {
                "label": "FC-3 (Bonus)  Random respawn angle",
                "syntax": "const, var, func_call",
                "code": [
                    "# in reset_ball(), after position reset:",
                    "# define ANGLE_MIN and ANGLE_MAX constants",
                    "# pick random: randf_range(ANGLE_MIN, ANGLE_MAX)",
                    "# ball_speed_x = cos(deg_to_rad(angle))",
                    "# ball_speed_y = sin(deg_to_rad(angle))",
                ],
            },
        ],
    },
    2: {
        "title": "Final Challenge — all 6 holes",
        "day": "DAY 2",
        "grid": "3x2",
        "fcs": [
            {
                "label": "FC-1  spawn_personality_ghosts()",
                "syntax": "for_range",
                "code": [
                    "# for i in range(PERSONALITY_COUNT):",
                    "#     spawn_one_personality(i)",
                ],
            },
            {
                "label": "FC-2  step_all_personality_ghosts()",
                "syntax": "for_in",
                "code": [
                    "# for ghost in ghosts:",
                    "#     step_personality(ghost)",
                ],
            },
            {
                "label": "FC-3  count_ghosts_of(personality) -> int",
                "syntax": "for_in, var, return",
                "code": [
                    "# var count := 0",
                    "# for ghost in ghosts:",
                    "#     if personality matches:",
                    "#         count += 1",
                    "# return count",
                ],
            },
            {
                "label": "FC-4  reset_personality_ghosts()",
                "syntax": "for_in, var, plus_eq",
                "code": [
                    "# var i := 0",
                    "# for ghost in ghosts:",
                    "#     respawn_personality_ghost(ghost, i)",
                    "#     i += 1",
                ],
            },
            {
                "label": "FC-5  target_for(ghost) -> Vector2i",
                "syntax": "func_return, if",
                "code": [
                    '# var p = ghost.get_meta("personality")',
                    "# if p == BLINKY:",
                    "#     return blinky_target(ghost)",
                    "# elif p == PINKY: ...",
                    "# elif p == INKY:  ...",
                    "# elif p == CLYDE: ...",
                ],
            },
            {
                "label": "FC-6  is_clyde_close(ghost) -> bool",
                "syntax": "func_return",
                "code": [
                    "# var d := distance_to_player(ghost)",
                    "# return d < 8.0",
                ],
            },
        ],
    },
    3: {
        "title": "Final Challenge — all 4 holes",
        "day": "DAY 3",
        "grid": "2x2",
        "fcs": [
            {
                "label": "FC-1  endless_tick(delta) body",
                "syntax": "for_in",
                "code": [
                    "# for e in main.enemies:",
                    "#     endless_buff(e, delta)",
                    "# for t in main.towers:",
                    "#     buff_tower(t, delta)",
                ],
            },
            {
                "label": "FC-2  buff_all(enemy_list, delta)",
                "syntax": "for_in, func_param",
                "code": [
                    "# for e in enemy_list:",
                    "#     endless_buff(e, delta)",
                ],
            },
            {
                "label": "FC-3  get_in_radius(pos, r) -> Array",
                "syntax": "list_init, for_in, list_append, return",
                "code": [
                    "# var result: Array = []",
                    "# for e in main.enemies:",
                    "#     if pos.distance_to(e.position) <= r:",
                    "#         result.append(e)",
                    "# return result",
                ],
            },
            {
                "label": "FC-4  check_wave_end()",
                "syntax": "if, list_size, plus_eq, func_call",
                "code": [
                    "# if wave_in_progress and enemies empty:",
                    "#     wave_in_progress = false",
                    "#     wave_index += 1",
                    "#     if end of WAVES: you_win()",
                    "#     else: start_next_wave()",
                ],
            },
        ],
    },
    4: {
        "title": "Final Challenge — all 3 holes",
        "day": "DAY 4",
        "grid": "3x1",
        "fcs": [
            {
                "label": "FC-1  CUSTOM_CHARACTER dict",
                "syntax": "var",
                "code": [
                    "# Fill all 11 keys:",
                    '# "display_name": "YourName"',
                    '# "sprite": "tile_XXXX.png"',
                    "# \"tint\": Color(r, g, b)",
                    '# "walk_speed", "jump_impulse"',
                    '# "attack_type": "custom"',
                    '# "attack_damage", "attack_cooldown"',
                    '# "attack_range", "projectile_speed"',
                    '# "projectile_gravity_scale"',
                ],
            },
            {
                "label": "FC-2  Register in CHARACTERS",
                "syntax": "func_call",
                "code": [
                    "# In main.gd _ready():",
                    '# CHARACTERS["custom"] = CUSTOM_CHARACTER',
                    '# Add "custom" to keys array',
                    "# Update title_label.text:",
                    '#     "5 = YourName"',
                ],
            },
            {
                "label": "FC-3  custom attack branch",
                "syntax": "match_stmt, if, func_call",
                "code": [
                    "# In player.gd attack():",
                    '# "custom":',
                    "#     opponent.take_damage(attack_damage)",
                    "#     # — or any of these —",
                    "#     # spawn_projectile()",
                    "#     # hp += 5",
                    "#     # anything you want",
                ],
            },
        ],
    },
}


# ── Helpers ──────────────────────────────────────────────────────────────────

def _solid_fill(shape, rgb: RGBColor):
    """Apply solid fill to a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _no_line(shape):
    """Remove shape border."""
    shape.line.fill.background()


def _add_textbox(slide, x, y, w, h):
    txBox = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def _para_run(tf, text, font_name, font_size_pt, bold=False, color=None,
              align=PP_ALIGN.LEFT, space_before=0, space_after=0, first=False):
    """Add a paragraph with one run to a text frame."""
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p


def _set_margins(tf, left=60000, right=60000, top=40000, bottom=40000):
    txBody = tf._txBody
    txBody.set("lIns", str(left))
    txBody.set("rIns", str(right))
    txBody.set("tIns", str(top))
    txBody.set("bIns", str(bottom))


def add_header(slide, day_label: str):
    """Add iCode header bar."""
    # Black bar
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(0), Emu(0), Emu(SW), Emu(HEADER_H)
    )
    _solid_fill(bar, BAR_BLACK)
    _no_line(bar)

    # Logo image (if available)
    logo_x = Emu(60000)
    logo_y = Emu(50000)
    logo_h = Emu(340000)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), logo_x, logo_y,
                                 height=logo_h)
    else:
        # Fallback text
        tb, tf = _add_textbox(slide, 60000, 60000, 900000, 340000)
        _para_run(tf, "iCode", FONT_HEAD, 20, bold=True, color=WHITE, first=True)
        _set_margins(tf, 30000, 30000, 20000, 20000)

    # Day label (right side, red)
    tb, tf = _add_textbox(slide, SW - 1200000, 100000, 1100000, 250000)
    _para_run(tf, day_label, FONT_HEAD, 18, bold=True, color=ICODE_RED,
              align=PP_ALIGN.RIGHT, first=True)
    _set_margins(tf, 30000, 80000, 20000, 20000)


def add_title(slide, title: str):
    """Add slide title below header."""
    tb, tf = _add_textbox(slide, 150000, HEADER_H + 30000, SW - 300000, TITLE_H)
    _para_run(tf, title, FONT_HEAD, 22, bold=True, color=TEXT_BLACK, first=True)
    _set_margins(tf, 60000, 60000, 20000, 20000)


def add_fc_box(slide, fc: dict, x: int, y: int, w: int, h: int):
    """Draw one FC box: label strip | syntax strip | code area."""
    PAD = 60000   # inner horizontal padding

    # Heights
    label_h = max(int(h * 0.14), 130000)
    syntax_h = max(int(h * 0.11), 100000)
    code_h   = h - label_h - syntax_h

    # ── Label strip (dark bg, white text) ───────────────────────────────────
    label_shape = slide.shapes.add_shape(
        1, Emu(x), Emu(y), Emu(w), Emu(label_h)
    )
    _solid_fill(label_shape, LABEL_BG)
    _no_line(label_shape)

    tb, tf = _add_textbox(slide, x, y, w, label_h)
    _para_run(tf, fc["label"], FONT_HEAD, 10, bold=True, color=WHITE, first=True)
    _set_margins(tf, PAD, PAD, 30000, 20000)

    # ── Syntax strip (slightly lighter bg, syntax key text) ──────────────────
    syn_y = y + label_h
    syn_shape = slide.shapes.add_shape(
        1, Emu(x), Emu(syn_y), Emu(w), Emu(syntax_h)
    )
    _solid_fill(syn_shape, SYNTAX_BG)
    _no_line(syn_shape)

    tb2, tf2 = _add_textbox(slide, x, syn_y, w, syntax_h)
    # "Syntax: " label in grey, then keys in blue
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r_lbl = p.add_run()
    r_lbl.text = "Syntax: "
    r_lbl.font.name = FONT_HEAD
    r_lbl.font.size = Pt(9)
    r_lbl.font.color.rgb = GREY_MID
    r_keys = p.add_run()
    r_keys.text = fc["syntax"]
    r_keys.font.name = FONT_MONO
    r_keys.font.size = Pt(9)
    r_keys.font.bold = True
    r_keys.font.color.rgb = SYNTAX_KEY
    _set_margins(tf2, PAD, PAD, 22000, 10000)

    # ── Code area (dark bg, green comment text, monospace) ───────────────────
    code_y = syn_y + syntax_h
    code_shape = slide.shapes.add_shape(
        1, Emu(x), Emu(code_y), Emu(w), Emu(code_h)
    )
    _solid_fill(code_shape, CODE_BG)
    _no_line(code_shape)

    tb3, tf3 = _add_textbox(slide, x, code_y, w, code_h)
    tf3.word_wrap = False
    lines = fc["code"]
    for i, line in enumerate(lines):
        _para_run(tf3, line, FONT_MONO, 9, color=SYNTAX_COMMENT,
                  space_before=0, space_after=0, first=(i == 0))
    _set_margins(tf3, PAD, PAD, 30000, 20000)


def build_fc_slide(day: int, data: dict, out_dir: Path):
    """Build and save a single FC slide for one day."""
    prs = Presentation()
    prs.slide_width = Emu(SW)
    prs.slide_height = Emu(SH)

    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)

    # White background
    bg = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(SW), Emu(SH))
    _solid_fill(bg, RGBColor(0xFF, 0xFF, 0xFF))
    _no_line(bg)

    add_header(slide, data["day"])
    add_title(slide, data["title"])

    fcs   = data["fcs"]
    grid  = data["grid"]
    n     = len(fcs)

    # ── Grid geometry ────────────────────────────────────────────────────────
    MARGIN = 150000
    COL_GAP = 100000
    ROW_GAP = 80000
    avail_w = SW - 2 * MARGIN
    avail_h = CONTENT_H

    if grid == "3x2":
        cols, rows = 3, 2
        col_w = (avail_w - (cols - 1) * COL_GAP) // cols
        row_h = (avail_h - (rows - 1) * ROW_GAP) // rows
        positions = [
            (MARGIN + c * (col_w + COL_GAP),
             CONTENT_TOP + r * (row_h + ROW_GAP))
            for r in range(rows) for c in range(cols)
        ]
    elif grid == "2x2":
        cols, rows = 2, 2
        col_w = (avail_w - COL_GAP) // 2
        row_h = (avail_h - ROW_GAP) // 2
        positions = [
            (MARGIN + c * (col_w + COL_GAP),
             CONTENT_TOP + r * (row_h + ROW_GAP))
            for r in range(rows) for c in range(cols)
        ]
    else:  # 3x1
        cols, rows = 3, 1
        col_w = (avail_w - (cols - 1) * COL_GAP) // cols
        row_h = avail_h
        positions = [
            (MARGIN + c * (col_w + COL_GAP), CONTENT_TOP)
            for c in range(cols)
        ]

    for i, fc in enumerate(fcs):
        if i >= len(positions):
            break
        px, py = positions[i]
        add_fc_box(slide, fc, px, py, col_w, row_h)

    out_path = out_dir / f"D{day}FC.pptx"
    prs.save(out_path)
    print(f"  saved {out_path}")


# ── Main ─────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    out_dir = Path(__file__).parent / "out"
    out_dir.mkdir(exist_ok=True)
    for day, data in FC_DATA.items():
        print(f"Building D{day}FC …")
        build_fc_slide(day, data, out_dir)
    print("Done.")
