"""L1-L8 slide layout functions.

Each layout function takes a content dict + day/page numbers and emits one slide
onto a given Presentation. Master frame (header + logo + day tab + page number)
applied first via master.apply_master().

All position/size constants live in theme.py — never hardcode here.

Layouts:
- L1 Title          — heading + optional subtitle, centered
- L2 Body           — heading + bullets/paragraphs
- L3 Side-by-Side   — heading + two columns (text/code/image left + right)
- L4 Image          — full-bleed image + caption strip
- L5 Table          — heading + header-row table
- L6 Code           — heading + centered monospace code block
- L7 Step           — screenshot + step badge + caption
- L8 Action         — top prose + LHS code + RHS screenshot with red overlay
"""

import re

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

from pygments import lex
from pygments.lexers import GDScriptLexer
from pygments.token import Token

try:
    from PIL import Image as _PILImage
except Exception:  # Pillow optional; aspect-fit degrades to box-fill
    _PILImage = None

import theme
from master import apply_master


# ============================================================
#  Inline markdown -> runs (**bold**, *italic*, `code`)
# ============================================================

_INLINE = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*|`([^`]+?)`")


def _emit_inline(p, text, *, font_name, size, bold=False, italic=False, color=None):
    """Add runs to paragraph `p`, converting **bold** / *italic* / `code` markdown
    into real formatting instead of printing the literal asterisks/backticks."""
    color = color or theme.TEXT_BLACK
    text = text or ""

    def add(t, b, i, mono):
        if not t:
            return
        r = p.add_run()
        r.text = t
        r.font.name = theme.FONT_MONO if mono else font_name
        r.font.size = size
        r.font.bold = bold or b
        r.font.italic = italic or i
        r.font.color.rgb = color

    pos = 0
    for m in _INLINE.finditer(text):
        if m.start() > pos:
            add(text[pos:m.start()], False, False, False)
        if m.group(1) is not None:
            add(m.group(1), True, False, False)
        elif m.group(2) is not None:
            add(m.group(2), False, True, False)
        else:
            add(m.group(3), False, False, True)
        pos = m.end()
    if pos < len(text):
        add(text[pos:], False, False, False)


# ============================================================
#  GDScript syntax highlighter — maps Pygments tokens to theme colors
# ============================================================

def _token_color(tok_type):
    """Map a Pygments token type to a theme RGBColor.

    GDScript lexer emits: Token.Keyword, Token.Keyword.Type, Token.Name.Function,
    Token.Name.Builtin, Token.Literal.String.*, Token.Literal.Number.*,
    Token.Comment.*, Token.Operator, Token.Punctuation, Token.Text.
    """
    if tok_type in Token.Comment:
        return theme.SYNTAX_COMMENT
    if tok_type in Token.Literal.String:
        return theme.SYNTAX_STRING
    if tok_type in Token.Literal.Number:
        return theme.SYNTAX_NUMBER
    if tok_type in Token.Keyword.Type:
        return theme.SYNTAX_TYPE
    if tok_type in Token.Keyword:
        # Pygments groups GDScript control-flow under Keyword; distinguish by string
        return theme.SYNTAX_KEYWORD
    if tok_type in Token.Name.Function:
        return theme.SYNTAX_FUNCTION
    if tok_type in Token.Name.Builtin:
        return theme.SYNTAX_FUNCTION
    if tok_type in Token.Name.Class:
        return theme.SYNTAX_TYPE
    return theme.CODE_TEXT


# ============================================================
#  Shared helpers
# ============================================================

def _new_slide(prs):
    """Add a blank slide (layout 6 = blank in default template)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _add_textbox(slide, left, top, width, height, text, *,
                 font_name=None, font_size=None, bold=False,
                 color=None, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    """Add a text box with one paragraph of text. Returns the shape."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = alignment
    _emit_inline(p, text,
                 font_name=font_name or theme.FONT_BODY,
                 size=font_size or theme.SIZE_BODY,
                 bold=bold,
                 color=color or theme.TEXT_BLACK)
    return box


def _add_bullets(slide, left, top, width, height, bullets, *,
                 font_size=None, color=None):
    """Add a textbox with one paragraph per bullet, hyphen-prefixed."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        prefix = p.add_run()
        prefix.text = "•  "
        prefix.font.name = theme.FONT_BODY
        prefix.font.size = font_size or theme.SIZE_BODY
        prefix.font.color.rgb = color or theme.TEXT_BLACK
        _emit_inline(p, bullet,
                     font_name=theme.FONT_BODY,
                     size=font_size or theme.SIZE_BODY,
                     color=color or theme.TEXT_BLACK)
        p.space_after = Pt(8)
    return box


def _add_code_block(slide, left, top, width, height, code_text,
                    bg_color=None, fg_color=None, font_size=None,
                    highlight=True):
    """Add a dark-bg code block with monospace text.

    `highlight=True` (default) runs the code through the Pygments GDScript
    lexer and applies per-token colors. `highlight=False` falls back to a
    single foreground color (useful for non-GDScript content like Python
    comparison panes).
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color or theme.CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)

    if not highlight:
        # Plain rendering — one run per line, single color
        for i, line in enumerate(code_text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line if line else " "
            run.font.name = theme.FONT_MONO
            run.font.size = font_size or theme.SIZE_CODE
            run.font.color.rgb = fg_color or theme.CODE_TEXT
        return shape

    # Pygments-highlighted rendering — one paragraph per line, one run per token
    size = font_size or theme.SIZE_CODE
    lines = code_text.split("\n")
    # Pre-tokenize the whole block once, then split tokens across lines
    tokens = list(lex(code_text, GDScriptLexer()))

    # Group tokens by line. A token can contain newlines — split on \n.
    lines_tokens = [[]]
    for tok_type, tok_text in tokens:
        if "\n" in tok_text:
            parts = tok_text.split("\n")
            for j, part in enumerate(parts):
                if part:
                    lines_tokens[-1].append((tok_type, part))
                if j < len(parts) - 1:
                    lines_tokens.append([])
        else:
            lines_tokens[-1].append((tok_type, tok_text))

    for i, line_tokens in enumerate(lines_tokens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if not line_tokens:
            # blank line — add a single space so paragraph height holds
            run = p.add_run()
            run.text = " "
            run.font.name = theme.FONT_MONO
            run.font.size = size
            run.font.color.rgb = theme.CODE_TEXT
            continue
        for tok_type, tok_text in line_tokens:
            run = p.add_run()
            run.text = tok_text
            run.font.name = theme.FONT_MONO
            run.font.size = size
            run.font.color.rgb = _token_color(tok_type)
    return shape


def _draw_overlay_box(slide, left, top, width, height):
    """Draw one red overlay rectangle (no fill, 4pt red outline) at exact geometry."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top),
                                  int(width), int(height))
    box.fill.background()
    box.line.color.rgb = theme.OVERLAY_RED
    box.line.width = Pt(4)
    box.shadow.inherit = False
    return box


def _add_placeholder_image(slide, left, top, width, height, image_path, caption=""):
    """Insert image (aspect-ratio preserved, centred in the box), or a
    dashed-border placeholder box if missing."""
    if image_path and Path(image_path).exists():
        w, h = int(width), int(height)
        if _PILImage is not None:
            try:
                iw, ih = _PILImage.open(str(image_path)).size
                scale = min(width / iw, height / ih)
                w, h = int(iw * scale), int(ih * scale)
            except Exception:
                w, h = int(width), int(height)
        off_l = left + (int(width) - w) // 2
        off_t = top + (int(height) - h) // 2
        return slide.shapes.add_picture(str(image_path), off_l, off_t,
                                          width=w, height=h)
    # placeholder
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    shape.line.color.rgb = theme.OVERLAY_RED
    shape.line.width = Pt(2)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ MISSING: {image_path.name if image_path else 'unspecified'} ]"
    run.font.name = theme.FONT_BODY
    run.font.size = Pt(14)
    run.font.color.rgb = theme.TEXT_MUTED
    return shape


# ============================================================
#  L1 — Title
# ============================================================

def l1_title(prs, *, day=None, page=None, heading, subtitle=None, bg_image=None):
    """Big centered heading + optional subtitle. Used for Day Title, Section Divider,
    Word Reveal, Day Closer, Personalization Beat Header."""
    slide = _new_slide(prs)
    apply_master(slide, day, page)

    if bg_image and Path(bg_image).exists():
        # Full-bleed background image below the header
        slide.shapes.add_picture(str(bg_image), Emu(0), theme.BODY_TOP,
                                  width=theme.SLIDE_WIDTH,
                                  height=theme.BODY_HEIGHT)

    # Heading centered vertically in the body region
    h_box = _add_textbox(
        slide,
        theme.BODY_LEFT_MARGIN, theme.BODY_TOP + Inches(1.4),
        theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN, Inches(2.0),
        heading,
        font_name=theme.FONT_HEADING,
        font_size=theme.SIZE_DAY_TITLE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    if subtitle:
        _add_textbox(
            slide,
            theme.BODY_LEFT_MARGIN, theme.BODY_TOP + Inches(3.6),
            theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN, Inches(1.0),
            subtitle,
            font_size=theme.SIZE_SUBHEADING,
            color=theme.TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )
    return slide


# ============================================================
#  L2 — Body (heading + bullets/paragraphs)
# ============================================================

def l2_body(prs, *, day=None, page=None, heading, bullets=None, paragraph=None, icon=None):
    """Headline + bullets or paragraph. Used for historical context, how-used,
    pieces-you'll-use, concept definitions, quiz Q+A."""
    slide = _new_slide(prs)
    apply_master(slide, day, page)

    _add_textbox(
        slide,
        theme.BODY_LEFT_MARGIN, theme.BODY_TOP + Inches(0.2),
        theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN, Inches(0.9),
        heading,
        font_name=theme.FONT_HEADING,
        font_size=theme.SIZE_HEADING,
        bold=True,
    )

    body_top = theme.BODY_TOP + Inches(1.4)
    body_width = theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN
    body_height = theme.SLIDE_HEIGHT - body_top - theme.BODY_BOTTOM_MARGIN

    if icon and Path(icon).exists():
        # Icon on the left, text on the right
        icon_size = Inches(1.5)
        slide.shapes.add_picture(str(icon), theme.BODY_LEFT_MARGIN,
                                   body_top, height=icon_size)
        body_left = theme.BODY_LEFT_MARGIN + icon_size + Inches(0.4)
        body_width -= (icon_size + Inches(0.4))
    else:
        body_left = theme.BODY_LEFT_MARGIN

    if bullets:
        _add_bullets(slide, body_left, body_top, body_width, body_height, bullets)
    elif paragraph:
        _add_textbox(slide, body_left, body_top, body_width, body_height, paragraph,
                     font_size=theme.SIZE_BODY)
    return slide


# ============================================================
#  L3 — Side-by-Side
# ============================================================

def l3_side_by_side(prs, *, day=None, page=None, heading=None,
                     left_label=None, left_text=None, left_code=None,
                     right_label=None, right_text=None, right_code=None):
    """Two-column layout. Used for GDScript-vs-Python, before/after,
    shape-in-code alongside metaphor caption."""
    slide = _new_slide(prs)
    apply_master(slide, day, page)

    if heading:
        _add_textbox(
            slide,
            theme.BODY_LEFT_MARGIN, theme.BODY_TOP + Inches(0.2),
            theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN, Inches(0.8),
            heading,
            font_name=theme.FONT_HEADING,
            font_size=theme.SIZE_HEADING,
            bold=True,
        )
        col_top = theme.BODY_TOP + Inches(1.2)
    else:
        col_top = theme.BODY_TOP + Inches(0.3)

    gutter = Inches(0.4)
    body_width = theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN
    col_width = (body_width - gutter) / 2
    col_height = theme.SLIDE_HEIGHT - col_top - theme.BODY_BOTTOM_MARGIN

    def _render_col(left, label, text, code):
        if label:
            _add_textbox(slide, left, col_top, col_width, Inches(0.5), label,
                         font_size=theme.SIZE_SUBHEADING, bold=True,
                         color=theme.ICODE_RED)
            content_top = col_top + Inches(0.7)
        else:
            content_top = col_top
        content_height = col_height - (content_top - col_top)
        if code:
            _add_code_block(slide, left, content_top, col_width, content_height, code)
        elif text:
            _add_textbox(slide, left, content_top, col_width, content_height, text,
                         font_size=theme.SIZE_BODY)

    _render_col(theme.BODY_LEFT_MARGIN, left_label, left_text, left_code)
    _render_col(theme.BODY_LEFT_MARGIN + col_width + gutter,
                right_label, right_text, right_code)
    return slide


# ============================================================
#  L4 — Image (full-bleed + caption)
# ============================================================

def l4_image(prs, *, day=None, page=None, image_path, caption=None):
    """Full-bleed image + caption strip. Used for metaphor hooks, after-works
    payoff, diagrams."""
    slide = _new_slide(prs)
    apply_master(slide, day, page)

    caption_height = Inches(0.8) if caption else Inches(0)
    image_top = theme.BODY_TOP + Inches(0.2)
    image_height = theme.SLIDE_HEIGHT - image_top - caption_height - theme.BODY_BOTTOM_MARGIN

    _add_placeholder_image(
        slide,
        theme.BODY_LEFT_MARGIN, image_top,
        theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN, image_height,
        image_path,
    )

    if caption:
        _add_textbox(
            slide,
            theme.BODY_LEFT_MARGIN, image_top + image_height + Inches(0.1),
            theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN, caption_height,
            caption,
            font_size=theme.SIZE_CAPTION,
            color=theme.TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    return slide


# ============================================================
#  L5 — Table
# ============================================================

def l5_table(prs, *, day=None, page=None, heading=None, header_row, data_rows):
    """Heading + table with header row + N data rows."""
    slide = _new_slide(prs)
    apply_master(slide, day, page)

    if heading:
        _add_textbox(
            slide,
            theme.BODY_LEFT_MARGIN, theme.BODY_TOP + Inches(0.2),
            theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN, Inches(0.8),
            heading,
            font_name=theme.FONT_HEADING,
            font_size=theme.SIZE_HEADING,
            bold=True,
        )
        table_top = theme.BODY_TOP + Inches(1.2)
    else:
        table_top = theme.BODY_TOP + Inches(0.3)

    body_width = theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN
    n_rows = len(data_rows) + 1
    n_cols = len(header_row)
    table_height_max = theme.SLIDE_HEIGHT - table_top - theme.BODY_BOTTOM_MARGIN
    row_height = min(Inches(0.5), Emu(int(table_height_max / n_rows)))

    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          theme.BODY_LEFT_MARGIN, table_top,
                                          body_width, row_height * n_rows)
    table = table_shape.table

    for col_i, header in enumerate(header_row):
        cell = table.cell(0, col_i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.GREY_DARK
        tf = cell.text_frame
        tf.text = header
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = theme.FONT_HEADING
                r.font.size = theme.SIZE_BODY_SMALL
                r.font.bold = True
                r.font.color.rgb = theme.BG_WHITE

    for row_i, row in enumerate(data_rows, start=1):
        for col_i, val in enumerate(row):
            cell = table.cell(row_i, col_i)
            tf = cell.text_frame
            tf.text = str(val)
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name = theme.FONT_BODY
                    r.font.size = theme.SIZE_BODY_SMALL
                    r.font.color.rgb = theme.TEXT_BLACK
    return slide


# ============================================================
#  L6 — Code (centered monospace)
# ============================================================

def l6_code(prs, *, day=None, page=None, heading=None, code, caption=None):
    """Heading + centered monospace code block. Used for standalone board examples,
    scene tree, constants."""
    slide = _new_slide(prs)
    apply_master(slide, day, page)

    body_width = theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN

    if heading:
        _add_textbox(
            slide,
            theme.BODY_LEFT_MARGIN, theme.BODY_TOP + Inches(0.2),
            body_width, Inches(0.8),
            heading,
            font_name=theme.FONT_HEADING,
            font_size=theme.SIZE_HEADING,
            bold=True,
        )
        code_top = theme.BODY_TOP + Inches(1.2)
    else:
        code_top = theme.BODY_TOP + Inches(0.5)

    caption_height = Inches(0.6) if caption else Inches(0)
    code_height = theme.SLIDE_HEIGHT - code_top - caption_height - theme.BODY_BOTTOM_MARGIN

    code_width = body_width * 0.8
    code_left = (theme.SLIDE_WIDTH - code_width) / 2

    _add_code_block(slide, code_left, code_top, code_width, code_height, code)

    if caption:
        _add_textbox(
            slide,
            theme.BODY_LEFT_MARGIN, code_top + code_height + Inches(0.1),
            body_width, caption_height,
            caption,
            font_size=theme.SIZE_CAPTION,
            color=theme.TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )
    return slide


# ============================================================
#  L7 — Step (screenshot + step badge + caption)
# ============================================================

def l7_step(prs, *, day=None, page=None, step_label=None, screenshot=None,
             caption=None, red_overlay=False, overlays=None):
    """Screenshot ~60% + step badge + caption ~40%. Used for all walks,
    where-in-game screenshots, personalization steps, export."""
    slide = _new_slide(prs)
    apply_master(slide, day, page)

    body_width = theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN
    body_top = theme.BODY_TOP + Inches(0.3)
    body_height = theme.SLIDE_HEIGHT - body_top - theme.BODY_BOTTOM_MARGIN

    # Step badge top-left
    if step_label:
        badge_size = Inches(0.9)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         theme.BODY_LEFT_MARGIN, body_top,
                                         badge_size, badge_size)
        badge.fill.solid()
        badge.fill.fore_color.rgb = theme.ICODE_RED
        badge.line.fill.background()
        tf = badge.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step_label
        run.font.name = theme.FONT_HEADING
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = theme.BG_WHITE

    # No screenshot: render the caption full-width as a text slide (no placeholder box).
    if not screenshot:
        if caption:
            _add_textbox(slide, theme.BODY_LEFT_MARGIN, body_top, body_width, body_height,
                         caption, font_size=theme.SIZE_BODY, anchor=MSO_ANCHOR.MIDDLE)
        return slide

    # Screenshot left ~60%, caption right ~40%
    img_left = theme.BODY_LEFT_MARGIN + Inches(1.1) if step_label else theme.BODY_LEFT_MARGIN
    img_width = body_width * 0.60
    img_height = body_height - Inches(0.2)
    img_shape = _add_placeholder_image(slide, img_left, body_top, img_width, img_height,
                                          screenshot)

    if overlays:
        # Restore hand-positioned overlay boxes at their exact saved geometry
        for o in overlays:
            _draw_overlay_box(slide, o["left"], o["top"], o["width"], o["height"])
    elif red_overlay:
        # Default red rectangle overlay on the screenshot (user drags + resizes)
        overlay_width = img_width * 0.4
        overlay_height = img_height * 0.2
        overlay_left = img_left + (img_width - overlay_width) / 2
        overlay_top = body_top + (img_height - overlay_height) / 2
        _draw_overlay_box(slide, overlay_left, overlay_top, overlay_width, overlay_height)

    if caption:
        caption_left = img_left + img_width + Inches(0.3)
        caption_width = body_width - (caption_left - theme.BODY_LEFT_MARGIN)
        _add_textbox(slide, caption_left, body_top, caption_width, img_height, caption,
                     font_size=theme.SIZE_BODY,
                     anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ============================================================
#  L8 — Action (top prose + LHS code + RHS screenshot with overlay)
# ============================================================

def l8_action(prs, *, day=None, page=None, prose, lhs_code, rhs_screenshot,
               two_tone=False, overlays=None):
    """The per-chunk Action slide. Top = R6 prose, LHS = board example code,
    RHS = Godot screenshot with red overlay marking kid #@todo region.
    two_tone=True adds a gray overlay underneath the red for R5 partial holes."""
    slide = _new_slide(prs)
    apply_master(slide, day, page)

    body_width = theme.SLIDE_WIDTH - theme.BODY_LEFT_MARGIN - theme.BODY_RIGHT_MARGIN
    body_top = theme.BODY_TOP + Inches(0.2)

    # TOP — prose instruction
    prose_height = Inches(1.4)
    prose_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          theme.BODY_LEFT_MARGIN, body_top,
                                          body_width, prose_height)
    prose_shape.fill.solid()
    prose_shape.fill.fore_color.rgb = theme.GREY_LIGHT
    prose_shape.line.fill.background()
    tf = prose_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _emit_inline(p, prose, font_name=theme.FONT_BODY, size=theme.SIZE_BODY,
                 bold=True, color=theme.TEXT_BLACK)

    # LHS — board example code block
    lhs_top = body_top + prose_height + Inches(0.3)
    lhs_height = theme.SLIDE_HEIGHT - lhs_top - theme.BODY_BOTTOM_MARGIN
    gutter = Inches(0.3)

    # No screenshot (e.g. Final-Challenge action slides whose answer is shown
    # only as a board example, never as a Godot capture): render the board
    # example full-width and skip the RHS column + overlay entirely.
    if not rhs_screenshot:
        _add_textbox(slide, theme.BODY_LEFT_MARGIN, lhs_top, body_width, Inches(0.35),
                     "Pattern (board example):",
                     font_size=theme.SIZE_BODY_SMALL, bold=True, color=theme.ICODE_RED)
        _add_code_block(slide,
                        theme.BODY_LEFT_MARGIN, lhs_top + Inches(0.4),
                        body_width, lhs_height - Inches(0.4),
                        lhs_code)
        return slide

    col_width = (body_width - gutter) / 2

    # Small label above LHS
    _add_textbox(slide, theme.BODY_LEFT_MARGIN, lhs_top, col_width, Inches(0.35),
                 "Pattern (board example):",
                 font_size=theme.SIZE_BODY_SMALL, bold=True, color=theme.ICODE_RED)
    _add_code_block(slide,
                    theme.BODY_LEFT_MARGIN, lhs_top + Inches(0.4),
                    col_width, lhs_height - Inches(0.4),
                    lhs_code)

    # RHS — Godot screenshot
    rhs_left = theme.BODY_LEFT_MARGIN + col_width + gutter
    _add_textbox(slide, rhs_left, lhs_top, col_width, Inches(0.35),
                 "Where in your code (Godot):",
                 font_size=theme.SIZE_BODY_SMALL, bold=True, color=theme.ICODE_RED)
    img_top = lhs_top + Inches(0.4)
    img_height = lhs_height - Inches(0.4)
    _add_placeholder_image(slide, rhs_left, img_top, col_width, img_height,
                            rhs_screenshot)

    # Two-tone gray overlay (R5 partial — drawn under the red)
    if two_tone:
        gray_overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                rhs_left + Inches(0.2), img_top + Inches(0.3),
                                                col_width - Inches(0.4), img_height * 0.4)
        gray_overlay.fill.solid()
        gray_overlay.fill.fore_color.rgb = theme.OVERLAY_GRAY
        gray_overlay.fill.transparency = 0.7
        gray_overlay.line.fill.background()

    if overlays:
        # Restore hand-positioned overlay boxes at their exact saved geometry
        for o in overlays:
            _draw_overlay_box(slide, o["left"], o["top"], o["width"], o["height"])
    else:
        # Default red overlay (user drags + resizes after build)
        red_width = col_width * 0.7
        red_height = img_height * 0.25
        red_left = rhs_left + (col_width - red_width) / 2
        red_top = img_top + (img_height - red_height) / 2
        _draw_overlay_box(slide, red_left, red_top, red_width, red_height)

    return slide
